import os
from enum import Enum
from marko.renderer import Renderer
import marko.block
import re
import pptx
from pptx.util import Pt


class PPTXRenderer(Renderer):

    def __init__(self):
        self.template_filename = None
        self.default_layout = 1  # for default template
        self.folder = None
        self.indent = None
        # for presentation content
        self.pres = None
        self.slide = None
        self.text_frame = None
        self.paragraph = None
        self.use_first_para = False
        self.run = None
        self.use_first_run = False

    def setup(self, template_filename: str = None, default_layout: int = None, folder: str = None) -> None:
        self.template_filename = template_filename
        self.default_layout = default_layout
        self.folder = folder

    def get_next_pres(self):
        if self.template_filename:
            self.pres = pptx.Presentation(self.template_filename)
        else:
            self.pres = pptx.Presentation()
        self.slide = None
        self.text_frame = None
        self.on_new_slide()

    def get_next_slide(self):
        title_slide_layout = self.pres.slide_layouts[self.default_layout]
        self.slide = self.pres.slides.add_slide(title_slide_layout)
        self.text_frame = self.slide.placeholders[1].text_frame
        self.indent = 0
        self.on_new_slide()

    def on_new_slide(self):
        self.paragraph = None
        self.use_first_para = False
        self.run = None
        self.use_first_run = False

    def get_next_paragraph(self):
        # Special case: if this is first paragraph, then no need to create another one
        # PowerPoint may create some content (paragraph and run)
        if not self.use_first_para:
            self.use_first_para = True
            self.paragraph = self.text_frame.paragraphs[0]
        else:
            self.paragraph = self.text_frame.add_paragraph()
        self.paragraph.level = self.indent-1 if self.indent > 0 else 0

        if len(self.paragraph.runs) > 0:
            self.run = self.paragraph.runs[0]
        else:
            self.run = self.paragraph.add_run()
        self.use_first_run = False

    def get_next_run(self):
        if not self.use_first_run:
            self.use_first_run = True
            self.run = self.paragraph.runs[0]
        else:
            self.run = self.paragraph.add_run()

    # Rendering functions

    def render_children_helper(self, element):
        children = getattr(element, "children", None)
        for child in children:
            self.render(child)

    def render_children_helper_str(self, element) -> str:
        children = getattr(element, "children", None)
        if isinstance(children, str):
            return children
        assert len(children) == 1
        if isinstance(children[0], marko.inline.RawText):
            return self.render_children_helper_str(children[0])
        raise RuntimeError('No string value')

    def render_children(self, element):
        # For debug purpose: not to miss something
        raise TypeError(f'Unknown type {element.__class__}')

    def render_document(self, element):
        self.get_next_pres()
        self.render_children_helper(element)

    def render_thematic_break(self, element):
        return

    def render_setext_heading(self, element):
        # NOTE: don't use it now
        # return self.render_heading(element)
        return

    def render_heading(self, element):
        self.get_next_slide()
        assert len(element.children) == 1
        self.slide.shapes.title.text = element.children[0].children

    def render_blank_line(self, element):
        # Do nothing
        pass

    def render_line_break(self, element):
        self.get_next_paragraph()

    def render_raw_text(self, element):
        self.get_next_run()
        self.run.text = element.children

    def render_list(self, element):
        self.indent += 1
        self.render_children_helper(element)
        self.indent -= 1

    def render_list_item(self, element):
        self.render_children_helper(element)

    def render_paragraph(self, element):
        self.get_next_paragraph()
        self.render_children_helper(element)

    def render_fenced_code(self, element):
        left = top = Pt(100)
        shape = self.slide.shapes.add_textbox(left, top, self.pres.slide_width, self.pres.slide_height)
        shape.text_frame.text = self.render_children_helper_str(element)

    def render_code_span(self, element):
        self.get_next_run()
        self.run.text = self.render_children_helper_str(element)
        assert self.run.text is not None
        font = self.run.font
        font.name = 'Consolas'

    def render_emphasis(self, element):
        self.get_next_run()
        self.run.text = self.render_children_helper_str(element)
        assert self.run.text is not None
        font = self.run.font
        font.italic = True

    def render_strong_emphasis(self, element):
        self.get_next_run()
        self.run.text = self.render_children_helper_str(element)
        assert self.run.text is not None
        font = self.run.font
        font.bold = True

    def render_html_block(self, element):
        match = re.search(r'<!--(.*?)-->\n?', element.children, flags=re.S)
        if match:
            self.slide.notes_slide.notes_text_frame.text += match.groups("1")[0].strip() + '\n'
            return
        raise RuntimeError(f'Unexpected html block: {element.children}')

    def render_image(self, element):
        left = top = Pt(100)
        filename = os.path.join(self.folder if self.folder else '', element.dest)
        pic = self.slide.shapes.add_picture(filename, left, top)
        if pic.width > self.pres.slide_width:
            k = pic.width / self.pres.slide_width
            pic.width = int(pic.width / k)
            pic.height = int(pic.height / k)
        elif pic.height > self.pres.slide_height:
            k = pic.height / self.pres.slide_height
            pic.width = int(pic.width / k)
            pic.height = int(pic.height / k)

    def render_link(self, element):
        self.get_next_run()
        self.run.text = self.render_children_helper_str(element)
        self.run.hyperlink.address = element.dest

    def render_auto_link(self, element):
        self.render_link(element)

    def render_link_ref_def(self, element):
        return
